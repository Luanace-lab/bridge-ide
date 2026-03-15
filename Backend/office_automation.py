"""
office_automation.py — Office Document Automation

Provides Excel, PowerPoint, and PDF manipulation for agents.
Wraps openpyxl, python-pptx, and pypdf with a unified interface.

Architecture Reference: R3_RealWorld_Capabilities.md
Phase: B — Capabilities

Features:
  - Excel: read/write cells, sheets, formulas, create workbooks
  - PowerPoint: create presentations, add slides, text, tables
  - PDF: extract text, page count, merge, split
  - Graceful degradation when libraries are not installed
  - Operation logging for audit trail

Design:
  - Optional dependencies — each library independently optional
  - Read operations always safe
  - Write operations logged
  - Thread-safe
"""

from __future__ import annotations

import os
import threading
import time
from dataclasses import dataclass, field
from enum import Enum
from typing import Any

# Optional imports — each library independently optional
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    openpyxl = None  # type: ignore[assignment]
    HAS_OPENPYXL = False

try:
    import pptx  # python-pptx
    HAS_PPTX = True
except ImportError:
    pptx = None  # type: ignore[assignment]
    HAS_PPTX = False

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    pypdf = None  # type: ignore[assignment]
    HAS_PYPDF = False


# ---------------------------------------------------------------------------
# Enums
# ---------------------------------------------------------------------------

class DocType(Enum):
    """Document type."""

    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    PDF = "pdf"


class OperationType(Enum):
    """Operation type for logging."""

    READ = "read"
    WRITE = "write"
    CREATE = "create"


# ---------------------------------------------------------------------------
# Data Classes
# ---------------------------------------------------------------------------

@dataclass
class OperationLog:
    """Audit log entry for document operations."""

    timestamp: float
    doc_type: str
    operation: str
    file_path: str
    agent_id: str = ""
    details: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "timestamp": self.timestamp,
            "doc_type": self.doc_type,
            "operation": self.operation,
            "file_path": self.file_path,
            "agent_id": self.agent_id,
            "details": self.details,
        }


@dataclass
class OperationResult:
    """Result of a document operation."""

    success: bool
    data: Any = None
    error: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "success": self.success,
            "data": self.data if not callable(self.data) else str(self.data),
            "error": self.error,
        }


@dataclass
class SheetData:
    """Data from an Excel sheet."""

    name: str
    rows: list[list[Any]] = field(default_factory=list)
    headers: list[str] = field(default_factory=list)
    row_count: int = 0
    col_count: int = 0

    def to_dict(self) -> dict[str, Any]:
        return {
            "name": self.name,
            "headers": self.headers,
            "row_count": self.row_count,
            "col_count": self.col_count,
            "rows": self.rows,
        }


# ---------------------------------------------------------------------------
# Office Automation Client
# ---------------------------------------------------------------------------

class OfficeClient:
    """Office document automation client.

    Provides safe access to Excel, PowerPoint, and PDF
    operations with audit logging.
    """

    def __init__(self) -> None:
        self._operation_log: list[OperationLog] = []
        self._lock = threading.Lock()

    # -------------------------------------------------------------------
    # Availability
    # -------------------------------------------------------------------

    @staticmethod
    def available_formats() -> dict[str, bool]:
        """Check which document formats are supported."""
        return {
            "excel": HAS_OPENPYXL,
            "powerpoint": HAS_PPTX,
            "pdf": HAS_PYPDF,
        }

    # -------------------------------------------------------------------
    # Excel Operations
    # -------------------------------------------------------------------

    def excel_read(
        self,
        file_path: str,
        sheet_name: str = "",
        agent_id: str = "",
    ) -> OperationResult:
        """Read an Excel file.

        Args:
            file_path: Path to .xlsx file.
            sheet_name: Sheet to read. Empty = active sheet.
            agent_id: Agent making the call.

        Returns:
            OperationResult with SheetData.
        """
        if not HAS_OPENPYXL:
            return OperationResult(success=False, error="openpyxl not installed")

        if not os.path.exists(file_path):
            return OperationResult(success=False, error=f"File not found: {file_path}")

        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

            rows = []
            headers = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                row_data = [cell for cell in row]
                if i == 0:
                    headers = [str(c) if c is not None else "" for c in row_data]
                rows.append(row_data)

            sheet_data = SheetData(
                name=ws.title or "",
                rows=rows,
                headers=headers,
                row_count=len(rows),
                col_count=len(headers),
            )
            wb.close()

            self._log(DocType.EXCEL.value, OperationType.READ.value,
                      file_path, agent_id, f"Sheet: {sheet_data.name}")
            return OperationResult(success=True, data=sheet_data)

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    def excel_write(
        self,
        file_path: str,
        data: list[list[Any]],
        sheet_name: str = "Sheet1",
        agent_id: str = "",
    ) -> OperationResult:
        """Write data to an Excel file.

        Args:
            file_path: Path for .xlsx output.
            data: 2D list of cell values. First row = headers.
            sheet_name: Sheet name.
            agent_id: Agent making the call.

        Returns:
            OperationResult.
        """
        if not HAS_OPENPYXL:
            return OperationResult(success=False, error="openpyxl not installed")

        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name

            for row in data:
                ws.append(row)

            wb.save(file_path)
            wb.close()

            self._log(DocType.EXCEL.value, OperationType.WRITE.value,
                      file_path, agent_id,
                      f"Sheet: {sheet_name}, Rows: {len(data)}")
            return OperationResult(success=True, data={"rows_written": len(data)})

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    def excel_get_sheets(
        self,
        file_path: str,
        agent_id: str = "",
    ) -> OperationResult:
        """Get sheet names from an Excel file.

        Args:
            file_path: Path to .xlsx file.
            agent_id: Agent making the call.

        Returns:
            OperationResult with list of sheet names.
        """
        if not HAS_OPENPYXL:
            return OperationResult(success=False, error="openpyxl not installed")

        if not os.path.exists(file_path):
            return OperationResult(success=False, error=f"File not found: {file_path}")

        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            names = wb.sheetnames
            wb.close()

            self._log(DocType.EXCEL.value, OperationType.READ.value,
                      file_path, agent_id, f"Sheets: {len(names)}")
            return OperationResult(success=True, data=names)

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    # -------------------------------------------------------------------
    # PowerPoint Operations
    # -------------------------------------------------------------------

    def pptx_create(
        self,
        file_path: str,
        slides: list[dict[str, Any]],
        agent_id: str = "",
    ) -> OperationResult:
        """Create a PowerPoint presentation.

        Args:
            file_path: Output .pptx path.
            slides: List of slide dicts with keys:
                - title: Slide title text
                - content: Slide body text (optional)
                - layout: Layout index (optional, default 1 = title+content)
            agent_id: Agent making the call.

        Returns:
            OperationResult.
        """
        if not HAS_PPTX:
            return OperationResult(success=False, error="python-pptx not installed")

        try:
            presentation = pptx.Presentation()

            for slide_data in slides:
                layout_idx = slide_data.get("layout", 1)
                if layout_idx >= len(presentation.slide_layouts):
                    layout_idx = 0
                layout = presentation.slide_layouts[layout_idx]
                slide = presentation.slides.add_slide(layout)

                if slide.placeholders:
                    if 0 in slide.placeholders:
                        slide.placeholders[0].text = slide_data.get("title", "")
                    if 1 in slide.placeholders and "content" in slide_data:
                        slide.placeholders[1].text = slide_data["content"]

            presentation.save(file_path)

            self._log(DocType.POWERPOINT.value, OperationType.CREATE.value,
                      file_path, agent_id, f"Slides: {len(slides)}")
            return OperationResult(
                success=True,
                data={"slides_created": len(slides)},
            )

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    def pptx_read(
        self,
        file_path: str,
        agent_id: str = "",
    ) -> OperationResult:
        """Read text content from a PowerPoint file.

        Args:
            file_path: Path to .pptx file.
            agent_id: Agent making the call.

        Returns:
            OperationResult with list of slide text dicts.
        """
        if not HAS_PPTX:
            return OperationResult(success=False, error="python-pptx not installed")

        if not os.path.exists(file_path):
            return OperationResult(success=False, error=f"File not found: {file_path}")

        try:
            presentation = pptx.Presentation(file_path)
            slides_data = []

            for i, slide in enumerate(presentation.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                slides_data.append({
                    "slide_number": i + 1,
                    "texts": texts,
                })

            self._log(DocType.POWERPOINT.value, OperationType.READ.value,
                      file_path, agent_id, f"Slides: {len(slides_data)}")
            return OperationResult(success=True, data=slides_data)

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    # -------------------------------------------------------------------
    # PDF Operations
    # -------------------------------------------------------------------

    def pdf_extract_text(
        self,
        file_path: str,
        pages: list[int] | None = None,
        agent_id: str = "",
    ) -> OperationResult:
        """Extract text from a PDF file.

        Args:
            file_path: Path to .pdf file.
            pages: Page numbers to extract (0-indexed). None = all.
            agent_id: Agent making the call.

        Returns:
            OperationResult with extracted text per page.
        """
        if not HAS_PYPDF:
            return OperationResult(success=False, error="pypdf not installed")

        if not os.path.exists(file_path):
            return OperationResult(success=False, error=f"File not found: {file_path}")

        try:
            reader = pypdf.PdfReader(file_path)
            total_pages = len(reader.pages)
            target_pages = pages if pages is not None else list(range(total_pages))

            result_pages = []
            for page_num in target_pages:
                if 0 <= page_num < total_pages:
                    text = reader.pages[page_num].extract_text() or ""
                    result_pages.append({
                        "page": page_num,
                        "text": text,
                    })

            self._log(DocType.PDF.value, OperationType.READ.value,
                      file_path, agent_id,
                      f"Pages: {len(result_pages)}/{total_pages}")
            return OperationResult(success=True, data={
                "total_pages": total_pages,
                "pages": result_pages,
            })

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    def pdf_page_count(
        self,
        file_path: str,
        agent_id: str = "",
    ) -> OperationResult:
        """Get page count of a PDF file.

        Args:
            file_path: Path to .pdf file.
            agent_id: Agent making the call.

        Returns:
            OperationResult with page count.
        """
        if not HAS_PYPDF:
            return OperationResult(success=False, error="pypdf not installed")

        if not os.path.exists(file_path):
            return OperationResult(success=False, error=f"File not found: {file_path}")

        try:
            reader = pypdf.PdfReader(file_path)
            count = len(reader.pages)

            self._log(DocType.PDF.value, OperationType.READ.value,
                      file_path, agent_id, f"Pages: {count}")
            return OperationResult(success=True, data={"page_count": count})

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    def pdf_merge(
        self,
        file_paths: list[str],
        output_path: str,
        agent_id: str = "",
    ) -> OperationResult:
        """Merge multiple PDF files.

        Args:
            file_paths: List of PDF file paths to merge.
            output_path: Output file path.
            agent_id: Agent making the call.

        Returns:
            OperationResult.
        """
        if not HAS_PYPDF:
            return OperationResult(success=False, error="pypdf not installed")

        for fp in file_paths:
            if not os.path.exists(fp):
                return OperationResult(success=False, error=f"File not found: {fp}")

        try:
            merger = pypdf.PdfMerger()
            for fp in file_paths:
                merger.append(fp)
            merger.write(output_path)
            merger.close()

            self._log(DocType.PDF.value, OperationType.WRITE.value,
                      output_path, agent_id,
                      f"Merged {len(file_paths)} files")
            return OperationResult(
                success=True,
                data={"files_merged": len(file_paths)},
            )

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    def pdf_split(
        self,
        file_path: str,
        output_dir: str,
        agent_id: str = "",
    ) -> OperationResult:
        """Split a PDF into individual pages.

        Args:
            file_path: Path to PDF file.
            output_dir: Directory for output files.
            agent_id: Agent making the call.

        Returns:
            OperationResult with list of output files.
        """
        if not HAS_PYPDF:
            return OperationResult(success=False, error="pypdf not installed")

        if not os.path.exists(file_path):
            return OperationResult(success=False, error=f"File not found: {file_path}")

        try:
            os.makedirs(output_dir, exist_ok=True)
            reader = pypdf.PdfReader(file_path)
            output_files = []

            base_name = os.path.splitext(os.path.basename(file_path))[0]

            for i, page in enumerate(reader.pages):
                writer = pypdf.PdfWriter()
                writer.add_page(page)
                out_path = os.path.join(output_dir, f"{base_name}_page_{i + 1}.pdf")
                with open(out_path, "wb") as f:
                    writer.write(f)
                output_files.append(out_path)

            self._log(DocType.PDF.value, OperationType.WRITE.value,
                      file_path, agent_id,
                      f"Split into {len(output_files)} pages")
            return OperationResult(
                success=True,
                data={"pages": len(output_files), "files": output_files},
            )

        except Exception as e:
            return OperationResult(success=False, error=str(e))

    # -------------------------------------------------------------------
    # Operation Log
    # -------------------------------------------------------------------

    def get_operation_log(
        self,
        limit: int = 50,
        doc_type: str = "",
        agent_id: str = "",
    ) -> list[OperationLog]:
        """Get operation log entries.

        Args:
            limit: Maximum entries to return.
            doc_type: Filter by document type. Empty = all.
            agent_id: Filter by agent. Empty = all.

        Returns:
            List of OperationLog entries (newest first).
        """
        with self._lock:
            log = list(reversed(self._operation_log))
        if doc_type:
            log = [e for e in log if e.doc_type == doc_type]
        if agent_id:
            log = [e for e in log if e.agent_id == agent_id]
        return log[:limit]

    # -------------------------------------------------------------------
    # Status
    # -------------------------------------------------------------------

    def status(self) -> dict[str, Any]:
        """Return office automation status summary."""
        return {
            "available_formats": self.available_formats(),
            "total_operations": len(self._operation_log),
            "supported_extensions": {
                "excel": [".xlsx"],
                "powerpoint": [".pptx"],
                "pdf": [".pdf"],
            },
        }

    # -------------------------------------------------------------------
    # Internal
    # -------------------------------------------------------------------

    def _log(
        self,
        doc_type: str,
        operation: str,
        file_path: str,
        agent_id: str,
        details: str = "",
    ) -> None:
        """Log a document operation."""
        entry = OperationLog(
            timestamp=time.time(),
            doc_type=doc_type,
            operation=operation,
            file_path=file_path,
            agent_id=agent_id,
            details=details,
        )
        with self._lock:
            self._operation_log.append(entry)
